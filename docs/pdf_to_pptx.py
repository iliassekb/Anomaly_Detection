"""
Convert a PDF (compiled from Beamer/LaTeX) to PPTX.
Uses PyMuPDF — no Poppler or external binaries needed.

Install once:  pip install pymupdf python-pptx
Run:           python pdf_to_pptx.py presentation.pdf presentation.pptx
"""

import sys
import io
import os
import fitz                          # pip install pymupdf
from pptx import Presentation
from pptx.util import Inches

PDF_PATH  = "presentation.pdf"
PPTX_PATH = "presentation.pptx"
DPI       = 200                      # 150=faster, 200=good, 300=print quality


def convert(pdf_path, pptx_path, dpi=200):
    print(f"Opening {pdf_path} ...")
    doc = fitz.open(pdf_path)

    prs = Presentation()
    prs.slide_width  = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]        # fully blank layout

    zoom = dpi / 72                     # fitz default is 72 dpi
    mat  = fitz.Matrix(zoom, zoom)

    for i, page in enumerate(doc, 1):
        pix  = page.get_pixmap(matrix=mat, alpha=False)
        buf  = io.BytesIO(pix.tobytes("png"))
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height)
        print(f"  slide {i}/{len(doc)}", end="\r")

    doc.close()
    prs.save(pptx_path)
    print(f"\nDone — saved to {pptx_path}")


if __name__ == "__main__":
    pdf  = sys.argv[1] if len(sys.argv) > 1 else PDF_PATH
    pptx = sys.argv[2] if len(sys.argv) > 2 else PPTX_PATH
    if not os.path.exists(pdf):
        print(f"Error: '{pdf}' not found.")
        print("Put the compiled PDF in the same folder as this script.")
        sys.exit(1)
    convert(pdf, pptx, dpi=DPI)
